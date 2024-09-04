from pptx import Presentation

def test_add_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert slide is not None
