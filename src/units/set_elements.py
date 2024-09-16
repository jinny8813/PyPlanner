from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import calendar
from datetime import timedelta
from pptx.dml.color import RGBColor

from units.weekday import date_type
from set_attribute import user_choice
from set_attribute import orig_font_color_element

def set_date_element_grid(slide,date_value,left, top, width, height, font_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = date_value
    p.font.size = Pt(24)
    p.font.name = "Arial"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_element
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 92, 92)
    else:
        p.font.color.rgb = RGBColor(92, 92, 255)
    p.alignment = PP_ALIGN.CENTER

def set_date_element(slide,date_value,left, top, width, height, lunar_text=None, font_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = date_value
    p.font.size = Pt(12)
    p.font.name = "Arial"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_element
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 92, 92)
    else:
        p.font.color.rgb = RGBColor(92, 92, 255)
    p.alignment = PP_ALIGN.CENTER
    if lunar_text!= None:
        run = p.add_run()
        run.text = " " + lunar_text
        run.font.size = Pt(8)

def set_lunar_element(slide,lunar_text,left, top, width, height, font_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = lunar_text
    p.font.size = Pt(8)
    p.font.name = "Arial"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_element
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 92, 92)
    else:
        p.font.color.rgb = RGBColor(92, 92, 255)

def set_small_calendar(slide,month_first_day,left, top,week_list):   
    o_left = left
    width = height = Pt(20)
    today = month_first_day
    for e in week_list:
        if e == "S":
            set_date_element(slide,e,left, top, width, height,None,"red")
        else:
            set_date_element(slide,e,left, top, width, height)
        left += width
    top += height
    for j in range(calendar.monthrange(month_first_day.year,month_first_day.month)[1]):
        if j == 0:
            if week_list[0] == "M":
                left = o_left + width * month_first_day.weekday()
            else:
                left = o_left + width * (month_first_day.weekday()+1)
        if left > o_left + width *6:
            top += height
            left = o_left 
            if j == 0:
                top -= height
        if user_choice["language"] == "holiday":
            date_details = date_type(today)
            set_date_element(slide,today.strftime("%#d"),left, top, width, height,None,date_details[2])
        else:
            set_date_element(slide,today.strftime("%#d"),left, top, width, height)
        today += timedelta(days=1)
        left += width