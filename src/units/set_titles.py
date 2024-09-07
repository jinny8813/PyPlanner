from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from set_attribute import orig_font_color_section, orig_font_color_page_title

def set_section_title(slide, section_info):
    text_top = Pt(216)
    text_left = Pt(512)
    text_height = Pt(96)
    text_width = Pt(512)
    tb = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    p = tb.text_frame.paragraphs[0]
    p.text = section_info
    p.font.size = Pt(96)
    p.font.name = "Optima"
    p.font.color.rgb = orig_font_color_section
    p.alignment = PP_ALIGN.CENTER

def set_chapter_title(slide, chapter_info):
    left = width = height = Pt(118)
    subtitle_top = Pt(540)
    tb = slide.shapes.add_textbox(left, subtitle_top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = chapter_info[0]
    p.font.size = Pt(32)
    p.font.name = "Optima"
    p.font.color.rgb = orig_font_color_section

    title_top = Pt(578)
    tb = slide.shapes.add_textbox(left, title_top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = chapter_info[1]
    p.font.size = Pt(72)
    p.font.name = "Optima"
    p.font.color.rgb = orig_font_color_section

def set_page_title(slide, title, lunar_text=None, font_color=None):
    width = height = Pt(24)
    left = Pt(45)
    top = Pt(50)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.name = "Optima"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_page_title
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 0, 0)
    else:
        p.font.color.rgb = RGBColor(0, 0, 255)
    if lunar_text!= None:
        run = p.add_run()
        run.text = " " + lunar_text
        run.font.size = Pt(16)