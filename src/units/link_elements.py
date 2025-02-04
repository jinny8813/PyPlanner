from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

from set_attribute import orig_font_color_element,orig_font_color_section,orig_nav_bg_color,orig_nav_bg_font_color,orig_font_color_little

def link_date_element(prs,slide,date_value,left, top, width, height, lunar_text=None, link_to = None, font_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = date_value
    p.font.size = Pt(12)
    p.font.name = "Arial"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_element
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 92, 92)
    else:
        p.font.color.rgb = RGBColor(92, 92, 255)
    p.alignment = PP_ALIGN.CENTER
    if lunar_text!= None:
        run = p.add_run()
        run.text = " " + lunar_text
        run.font.size = Pt(8)
    if link_to!= None:
        tb.click_action.target_slide = prs.slides[link_to]

def link_date_element_grid(prs,slide,date_value,left, top, width, height, link_to = None,font_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = date_value
    p.font.size = Pt(24)
    p.font.name = "Arial"
    if font_color == None or font_color == "black":
        p.font.color.rgb = orig_font_color_element
    elif font_color == "red":
        p.font.color.rgb = RGBColor(255, 92, 92)
    else:
        p.font.color.rgb = RGBColor(92, 92, 255)
    p.alignment = PP_ALIGN.CENTER
    if link_to!= None:
        tb.click_action.target_slide = prs.slides[link_to]

def link_calendar_weeks(prs,slide,week_value,left, top, width, height, link_to):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = week_value
    p.font.size = Pt(10)
    p.font.name = "Arial"
    p.font.color.rgb = orig_font_color_element
    p.alignment = PP_ALIGN.CENTER
    tb.click_action.target_slide = prs.slides[link_to]

def link_top_nav(prs,slide,feature_value,left, top, width, height, link_to):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,left,top,width,height)
    dot.fill.solid()
    dot.fill.fore_color.rgb = orig_nav_bg_color
    dot.line.color.rgb = orig_nav_bg_color
    tb = slide.shapes.add_textbox(left, top+height*0.06, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = feature_value
    p.font.size = Pt(9)
    p.font.name = "Arial"
    p.font.color.rgb = orig_nav_bg_font_color
    p.alignment = PP_ALIGN.CENTER     
    tb.click_action.target_slide = prs.slides[link_to]

def link_main_nav(prs,slide,name,left, top, width, height,link_to):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.name = "Arial"
    p.font.color.rgb = orig_font_color_section
    p.alignment = PP_ALIGN.CENTER
    if link_to!= None:
        tb.click_action.target_slide = prs.slides[link_to]

def link_gn_apple(slide,left, top, width, height, link_to):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = "\uf179"
    p.font.size = Pt(7)
    p.font.name = "Font Awesome 6 Brands Regular"
    p.font.color.rgb = orig_font_color_little
    p.alignment = PP_ALIGN.CENTER     
    tb.click_action.hyperlink.address = link_to

def link_gn_google(slide,left, top, width, height, link_to):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = "\uf1a0"
    p.font.size = Pt(5.5)
    p.font.name = "Font Awesome 6 Brands Regular"
    p.font.color.rgb = orig_font_color_little
    p.alignment = PP_ALIGN.CENTER     
    tb.click_action.hyperlink.address = link_to