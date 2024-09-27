from pptx.util import Pt
from datetime import timedelta

from units.link_elements import link_gn_apple,link_gn_google
from set_attribute import orig_start_time_2

def link_gn_apple_timeline(prs,m_type_count,w_type_count):
    height = Pt(20)
    width = Pt(20)
    orig_top = top = Pt(200)
    orig_left = left = Pt(257)
    space = Pt(120)
    date = orig_start_time_2
    for slide in prs.slides:
        current_page = prs.slides.index(slide)
        if current_page>=6+13*m_type_count+54*w_type_count:
            break
        if current_page<6+13*m_type_count+54*(w_type_count-1)+1:
            continue
        left = orig_left
        for j in range(7):
            top = orig_top
            for i in range(24):
                dt = date.strftime("%Y-%m-%dT%H:%M:%S")
                link_gn_apple(slide,left, top+height*i, width, height,f"shortcuts://run-shortcut/?name=BluegaJournal&input=Blue01-{dt}")
                date += timedelta(hours=1)
            left += space

def link_gn_google_timeline(prs,m_type_count,w_type_count):
    height = Pt(20)
    width = Pt(20)
    orig_top = top = Pt(211.5)
    orig_left = left = Pt(257)
    space = Pt(120)
    date = orig_start_time_2 - timedelta(hours=8)
    for slide in prs.slides:
        current_page = prs.slides.index(slide)
        if current_page>=6+13*m_type_count+54*w_type_count:
            break
        if current_page<6+13*m_type_count+54*(w_type_count-1)+1:
            continue
        left = orig_left
        for j in range(7):
            top = orig_top
            for i in range(24):
                new_date = date + timedelta(hours=1)
                dt = date.strftime("%Y%m%dT%H%M%SZ")
                ndt = new_date.strftime("%Y%m%dT%H%M%SZ")
                link_gn_google(slide,left, top+height*i, width, height,f"https://calendar.google.com/calendar/u/0/r/eventedit?dates={dt}/{ndt}")
                date = new_date
            left += space