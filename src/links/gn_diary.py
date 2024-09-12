from pptx.util import Pt
from datetime import timedelta
from dateutil.relativedelta import relativedelta

from units.link_elements import link_gn_apple,link_gn_google
from set_attribute import orig_start_month, orig_start_time
from set_attribute import orig_selected_w_types,orig_selected_m_types

def link_gn_apple_diary(prs):
    height = Pt(20)
    width = Pt(20)
    top = Pt(200)
    left = Pt(137)
    date = orig_start_time
    year_count = (orig_start_month+relativedelta(years=1)-orig_start_month)
    for slide in prs.slides:
        current_page = prs.slides.index(slide)
        if current_page>=6+len(orig_selected_m_types)*13+len(orig_selected_w_types)*54+1+year_count.days:
            break
        if current_page<6+len(orig_selected_m_types)*13+len(orig_selected_w_types)*54+1:
            continue
        for i in range(24):
            dt = date.strftime("%Y-%m-%dT%H:%M:%S")
            link_gn_apple(slide,left, top+height*i, width, height,f"shortcuts://run-shortcut/?name=BluegaJournal&input=Blue01-{dt}")
            date += timedelta(hours=1)

def link_gn_google_diary(prs):
    height = Pt(20)
    width = Pt(20)
    top = Pt(211.5)
    left = Pt(137)
    date = orig_start_time - timedelta(hours=8)
    year_count = (orig_start_month+relativedelta(years=1)-orig_start_month)
    for slide in prs.slides:
        current_page = prs.slides.index(slide)
        if current_page>=6+len(orig_selected_m_types)*13+len(orig_selected_w_types)*54+1+year_count.days:
            break
        if current_page<6+len(orig_selected_m_types)*13+len(orig_selected_w_types)*54+1:
            continue
        for i in range(24):
            new_date = date + timedelta(hours=1)
            dt = date.strftime("%Y%m%dT%H%M%SZ")
            ndt = new_date.strftime("%Y%m%dT%H%M%SZ")
            link_gn_google(slide,left, top+height*i, width, height,f"shortcuts://run-shortcut/?name=BluegaJournal&input=Blue02-{dt}/{ndt}")
            date = new_date