from pptx.util import Pt
from dateutil.relativedelta import relativedelta
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from units.link_elements import link_main_nav
from set_attribute import orig_start_month
from set_attribute import orig_theme_colors
from set_attribute import orig_font_color_section

def link_down_nav(prs,my_index_list):
    main_nav_names = ["Index", "Journal", "Lifestyle", "Stickers"]
    main_nav_section = []
    main_nav_section.append([item for item in my_index_list if item[3] == "索引"][0])
    main_nav_section.append([item for item in my_index_list if item[3] == "手帳年曆"][0])
    main_nav_section.append([item for item in my_index_list if item[3] == "生活風格"][0])
    main_nav_section.append([item for item in my_index_list if item[3] == "貼紙合集"][0])
    index_life = my_index_list.index(main_nav_section[2])
    index_sticker = my_index_list.index(main_nav_section[3])
    this_month = orig_start_month
    for slide in prs.slides:
        current_page = prs.slides.index(slide)
        if current_page>len(prs.slides)-2:
            break
        if current_page<1:
            continue
        top = Pt(736)
        left = Pt(0)
        height = Pt(32)
        width = Pt(80)
        for i,row in enumerate(orig_theme_colors):
            rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(row[0], row[1], row[2])
            rectangle.line.color.rgb = RGBColor(row[0], row[1], row[2])
            if i<4:
                tb = slide.shapes.add_textbox(left, top+Pt(6), width, height-Pt(6))
                p = tb.text_frame.paragraphs[0]
                p.text = main_nav_names[i]
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.font.color.rgb = orig_font_color_section
                p.alignment = PP_ALIGN.CENTER
                tb.click_action.target_slide = prs.slides[main_nav_section[i][1]]
            left += width
            if i == 3:
                width = Pt(704)
        width = Pt(80)
        right = Pt(1024) - width
        if current_page<main_nav_section[1][1]:
            link_main_nav(prs,slide,main_nav_section[0][2],right, top+Pt(6), width, height-Pt(6),main_nav_section[0][1])
        elif current_page<main_nav_section[2][1] and current_page>=main_nav_section[1][1]:
            link_main_nav(prs,slide,main_nav_section[1][2],right, top+Pt(6), width, height-Pt(6),main_nav_section[1][1])
            space = Pt(32)
            right = Pt(1024) - width - space*12.5
            for j in range(12):
                link_main_nav(prs,slide,this_month.strftime("%#m"),right, top+Pt(6), space, height-Pt(6),6+j+1)
                this_month += relativedelta(months=1)
                right += space
        elif current_page<main_nav_section[3][1] and current_page>=main_nav_section[2][1]:
            link_main_nav(prs,slide,main_nav_section[2][2],right, top+Pt(6), width, height-Pt(6),main_nav_section[2][1])
            space = Pt(42)
            right = Pt(1024) - width - space*12 - Pt(16)
            short = ["Book","Movie","Music","Attr.","Portf.","Etc.","nb-A","nb-B","nb-C","nb-D","nb-E","Temp."]
            for j,item in enumerate(short):
                link_main_nav(prs,slide,item,right, top+Pt(6), space, height-Pt(6),my_index_list[index_life+j+1][1])
                right += space
        else:
            link_main_nav(prs,slide,main_nav_section[3][2],right, top+Pt(6), width, height-Pt(6),main_nav_section[3][1])
            space = Pt(42)
            right = Pt(1024) - width - space*6 - Pt(16)
            short = ["Cal.","Lbl.","Rect.","Sol.","Deco.","Oth."]
            for j,item in enumerate(short):
                link_main_nav(prs,slide,item,right, top+Pt(6), space, height-Pt(6),my_index_list[index_sticker+j+1][1])
                right += space