from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from set_attribute import orig_font_color_element
from set_attribute import user_choice

def link_index(prs,my_index_list):
    top = Pt(93)
    left = Pt(92)
    width = Pt(360)
    height = Pt(20)
    slide = prs.slides[1]
    for i,row in enumerate(my_index_list):
        if row[2] == "Lifestyle":
            left += Pt(480)
            top = Pt(93)
        if row[0] == "big":
            top += height
            tb = slide.shapes.add_textbox(left, top, width, height*2)
            p = tb.text_frame.paragraphs[0]
            p.text = row[2]
            p.font.size = Pt(18)
            p.font.name = "Arial"
            p.font.color.rgb = orig_font_color_element
            p.alignment = PP_ALIGN.CENTER
            if user_choice["language"] != "english":
                run = p.add_run()
                run.text = " " + row[3]
                run.font.size = Pt(12)
            tb.click_action.target_slide = prs.slides[row[1]]
            top += height*2
        elif row[0] == "small":
            tb = slide.shapes.add_textbox(left, top, width, height)
            p = tb.text_frame.paragraphs[0]
            p.text = " - " + row[2]
            p.font.size = Pt(14)
            p.font.name = "Arial"
            p.font.color.rgb = orig_font_color_element
            if user_choice["language"] != "english":
                run = p.add_run()
                run.text = " " + row[3]
                run.font.size = Pt(8)
            tb.click_action.target_slide = prs.slides[row[1]]
            top += height*1.25