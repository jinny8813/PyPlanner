from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from datetime import timedelta
from dateutil.relativedelta import relativedelta
import calendar

from units.set_titles import set_page_title
from units.set_elements import set_small_calendar,set_date_element
from set_attribute import orig_start_month,orig_week_list
from set_attribute import orig_font_color_element

def create_yaag_slides(prs, ppt_count):
    layout = prs.slide_layouts[21]
    slide = prs.slides.add_slide(layout)
    set_page_title(slide,"2025 Year at a Glance")
    this_month = orig_start_month
    left = Pt(42)
    width = height = Pt(20)
    for i in range(12):
        if i/6 >= 1:
            top = Pt(423)
        else:
            top = Pt(103)
        tb = slide.shapes.add_textbox(left+width*8*(i%6), top, width*7, height*2)
        p = tb.text_frame.paragraphs[0]
        p.text = this_month.strftime("%B")
        p.font.size = Pt(16)
        p.font.name = "Arial"
        p.font.color.rgb = orig_font_color_element
        p.alignment = PP_ALIGN.CENTER
        top += height*1.5
        set_small_calendar(slide,this_month,left+width*8*(i%6), top,orig_week_list)
        this_month += relativedelta(months=1)
    ppt_count += 1
    return ppt_count

def create_gc_slides(prs, ppt_count):
    layout = prs.slide_layouts[21]
    slide = prs.slides.add_slide(layout)
    set_page_title(slide,"2025 Gantt Chart")
    this_month = orig_start_month
    today = orig_start_month
    height = Pt(20)
    width = Pt(20)
    left = Pt(32)
    for i in range(12):
        top = Pt(83)
        set_date_element(slide,this_month.strftime("%b"),left, top, width*4, height)
        top += height
        for j in range(31):
            if j < calendar.monthrange(this_month.year,this_month.month)[1]:
                set_date_element(slide,today.strftime("%#d"),left, top, width, height)
                set_date_element(slide,today.strftime("%a")[0],left+width, top, width, height)
                today += timedelta(days=1)
            else:
                set_date_element(slide,"/",left, top, width, height)
                set_date_element(slide,"/",left+width, top, width, height)
            top += height
        left += width*4
        this_month += relativedelta(months=1)
    ppt_count += 1
    return ppt_count