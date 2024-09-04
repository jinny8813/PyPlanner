from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from units.set_titles import set_chapter_title,set_section_title

from set_attribute import orig_theme_colors

def create_calendar_slides(prs, my_index_list, ppt_count, chapter_titles, slide_titles):
    layout = prs.slide_layouts[19]
    slide = prs.slides.add_slide(layout)
    set_chapter_title(slide,chapter_titles)
    my_index_list.append(["small",ppt_count,slide_titles[0],slide_titles[1]])
    ppt_count += 1
    return ppt_count

def create_section_slides(prs, my_index_list, ppt_count, section_info, slide_titles,tcolor_count):
    layout = prs.slide_layouts[19]
    slide = prs.slides.add_slide(layout)
    top = Pt(148)
    left = Pt(0)
    height = Pt(620)
    width = Pt(1024)
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
    rectangle.fill.solid()
    rectangle.fill.fore_color.rgb = RGBColor(orig_theme_colors[tcolor_count][0], orig_theme_colors[tcolor_count][1], orig_theme_colors[tcolor_count][2])
    rectangle.line.color.rgb = RGBColor(orig_theme_colors[tcolor_count][0], orig_theme_colors[tcolor_count][1], orig_theme_colors[tcolor_count][2])
    set_section_title(slide,section_info)
    my_index_list.append(["big",ppt_count,slide_titles[0],slide_titles[1]])
    ppt_count += 1
    return ppt_count