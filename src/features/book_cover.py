from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from units.set_titles import set_chapter_title
from set_attribute import orig_theme_colors
from set_attribute import orig_font_color_section

def create_front_cover_slides(prs, my_index_list, ppt_count, chapter_titles, slide_titles):
    layout = prs.slide_layouts[19]
    slide = prs.slides.add_slide(layout)
    create_theme_colors_block(slide)
    set_chapter_title(slide,chapter_titles)
    my_index_list.append(["cover",ppt_count,slide_titles[0],slide_titles[1]])
    ppt_count += 1
    return ppt_count

def create_back_cover_slides(prs, my_index_list, ppt_count, slide_titles):
    layout = prs.slide_layouts[19]
    slide = prs.slides.add_slide(layout)
    create_theme_colors_block(slide)
    left = width = height = Pt(118)
    subtitle_top = Pt(540)
    tb = slide.shapes.add_textbox(left, subtitle_top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = "The Blueprint"
    p.font.size = Pt(32)
    p.font.name = "Optima"
    p.font.color.rgb = orig_font_color_section
    title_top = Pt(578)
    tb = slide.shapes.add_textbox(left, title_top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = "The dialogue between dreams and reality \nDesign by @bluega_journal 2024 \nEnjoy journaling"
    p.font.size = Pt(14)
    p.font.name = "Optima"
    p.font.color.rgb = orig_font_color_section
    tb.click_action.hyperlink.address = "https://www.instagram.com/bluega_journal/"
    my_index_list.append(["cover",ppt_count,slide_titles[0],slide_titles[1]])
    ppt_count += 1
    return ppt_count

def create_theme_colors_block(slide):
    top = Pt(148)
    left = Pt(0)
    height = Pt(620)
    width = Pt(864)
    for i,row in enumerate(orig_theme_colors):
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(row[0], row[1], row[2])
        rectangle.line.color.rgb = RGBColor(row[0], row[1], row[2])
        left += width
        width = Pt(40)